"""
文档翻译模块 v4.6.1
功能：Word/Excel/PPT 文档专业翻译，支持多引擎降级

v4.6.1 变更:
  - 🔒 翻译引擎改为显式 Opt-in：auto 模式仅使用 local-rule 和 pure-template，不读取 API Key
  - 🔒 LLM 翻译仅在用户 method 显式指定时调用，首次使用显示凭证读取范围警告
  - 🔒 移除外部 LLM 的自动检测与降级

v4.6.0 变更:
  - 🎯 多格式支持（Word/Excel/PPT）文档翻译
  - 🎯 保持原文档格式（字体/样式/表格/图表）
  - 🎯 术语表支持（JSON 格式，可配置）
  - 🎯 批量目录翻译
"""

import os
import sys
import json
import re
import hashlib
import time
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from datetime import datetime

sys.path.insert(0, str(Path(__file__).parent))

try:
    from wps_common import safe_path, get_hardware_info, with_retry
except ImportError:
    def safe_path(p): return Path(p)
    def get_hardware_info(): return {"cpu_cores": 4, "memory_gb": 8, "level": "medium"}
    def with_retry(f): return f


class TranslationEngine:
    """翻译引擎（显式 Opt-in 模式）"""
    
    def __init__(self, method: str = "auto"):
        self.method = method
        self._available = self._detect_available()
        self.glossary = self._load_glossary()
    
    def _detect_available(self) -> Dict[str, bool]:
        """检测可用的翻译引擎（仅本地引擎）"""
        return {
            "cn-llm-router": False,  # 仅在显式指定时检测
            "local-rule": True,  # 始终可用（术语表）
            "pure-template": True,  # 始终可用
        }
    
    def _detect_external_available(self, method: str) -> bool:
        """检测外部 LLM 引擎是否可用（仅在显式调用时执行）"""
        if method == "cn-llm-router":
            return bool(os.environ.get("LLM_ROUTER_API_KEY"))
        elif method == "external-llm":
            return bool(os.environ.get("OPENAI_API_KEY"))
        return False
    
    def _load_glossary(self) -> Dict[str, str]:
        """加载术语表"""
        glossary_path = Path(__file__).parent / "glossary.json"
        if glossary_path.exists():
            try:
                return json.loads(glossary_path.read_text(encoding="utf-8"))
            except Exception:
                pass
        # 默认术语表
        return {
            "figure": "图",
            "table": "表",
            "equation": "公式",
            "reference": "参考文献",
            "abstract": "摘要",
            "introduction": "引言",
            "conclusion": "结论",
            "appendix": "附录",
        }
    
    def get_best_method(self) -> str:
        """获取最佳翻译引擎（auto 模式仅使用本地引擎）"""
        if self.method != "auto":
            # 显式指定外部引擎时，检测是否可用
            if self.method in ("cn-llm-router", "external-llm"):
                if self._detect_external_available(self.method):
                    return self.method
                else:
                    print(f"[翻译] ⚠️ 外部引擎 '{self.method}' 不可用（API Key 未配置），降级到 local-rule", file=sys.stderr)
                    return "local-rule"
            return self.method
        
        # auto 模式：仅使用 local-rule，不读取 API Key
        return "local-rule"
    
    @with_retry
    def translate(self, text: str, src: str = "en", tgt: str = "zh") -> Dict:
        """翻译文本"""
        if not text.strip():
            return {"success": True, "text": "", "method": "none"}
        
        method = self.get_best_method()
        
        # 外部服务调用前显示警告
        if method in ("cn-llm-router", "external-llm"):
            self._warn_llm_usage(method)
        
        try:
            if method in ("cn-llm-router", "external-llm"):
                result = self._translate_llm(text, src, tgt)
            elif method == "local-rule":
                result = self._translate_rule(text, src, tgt)
            else:
                result = self._translate_template(text, src, tgt)
            
            return {"success": True, "method": method, **result}
        except Exception as e:
            if method != "pure-template":
                return self._translate_template(text, src, tgt)
            return {"success": False, "error": str(e)}
    
    def _warn_llm_usage(self, method: str):
        """LLM 翻译使用前的 Opt-in 警告"""
        base_url = os.environ.get("LLM_ROUTER_BASE_URL", 
                                 os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1"))
        api_key = os.environ.get("LLM_ROUTER_API_KEY", os.environ.get("OPENAI_API_KEY", ""))
        key_preview = key[:4] + "..." + key[-4:] if len(key) > 8 else "未配置"
        model = os.environ.get("LLM_MODEL", "gpt-4o-mini")
        
        print(f"""
⚠️  [外部 LLM 翻译服务调用警告]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📋 即将使用外部 LLM 服务进行文档翻译
🔑 API Key: {key_preview}
🌐 服务端点: {base_url}/chat/completions
📌 模型: {model}
📤 数据流向: 文档内容将发送至外部 LLM 服务
📌 用途: 技术文档翻译
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
API Key 将读取并发送至外部服务。
外部服务将在上方信息确认后开始调用。
""", file=sys.stderr)
    
    def _translate_llm(self, text: str, src: str, tgt: str) -> Dict:
        """LLM 翻译（仅显式调用）"""
        # 根据 method 确定使用哪个外部服务
        if self.method == "cn-llm-router":
            base_url = os.environ.get("LLM_ROUTER_BASE_URL", "https://api.openai.com/v1")
            api_key = os.environ.get("LLM_ROUTER_API_KEY", "")
        else:  # external-llm
            base_url = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")
            api_key = os.environ.get("OPENAI_API_KEY", "")
        
        model = os.environ.get("LLM_MODEL", "gpt-4o-mini")
        
        if not api_key:
            raise ValueError("未配置 API Key")
        
        import urllib.request
        import urllib.error
        
        # 构建术语表提示
        glossary_hint = ""
        if self.glossary:
            glossary_hint = "\n\n术语表（请严格遵循）：\n" + "\n".join(f"{k} → {v}" for k, v in list(self.glossary.items())[:20])
        
        prompt = f"""请将以下{src}文本翻译为{tgt}。要求：
1. 保持专业术语一致性
2. 保留原文格式（标题、列表、代码块等）
3. 技术文档风格，准确流畅
4. 仅输出翻译结果，不要解释

原文：
{text[:3000]}
"""
        
        payload = json.dumps({
            "model": model,
            "messages": [{"role": "user", "content": prompt + glossary_hint}],
            "temperature": 0.3,
            "max_tokens": 4000,
        }).encode("utf-8")
        
        req = urllib.request.Request(
            f"{base_url}/chat/completions",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}",
            }
        )
        
        try:
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                translated = result["choices"][0]["message"]["content"]
                return {"text": translated}
        except urllib.error.HTTPError as e:
            raise ValueError(f"API 错误: {e.code}")
    
    def _translate_rule(self, text: str, src: str, tgt: str) -> Dict:
        """规则翻译（术语表替换）"""
        if src == "en" and tgt == "zh":
            result = text
            for eng, chn in self.glossary.items():
                result = re.sub(r'\b' + re.escape(eng) + r'\b', chn, result, flags=re.IGNORECASE)
            return {"text": result, "partial": True}
        return {"text": text, "partial": False}
    
    def _translate_template(self, text: str, src: str, tgt: str) -> Dict:
        """模板翻译"""
        template = f"""[翻译模板 - 请手动填充]

源语言: {src}
目标语言: {tgt}
引擎: pure-template（无可用的翻译引擎）

请安装以下任一引擎：
  - cn-llm-router: 设置 LLM_ROUTER_API_KEY 环境变量
  - OpenAI: 设置 OPENAI_API_KEY 环境变量

原文（前 500 字符）：
{text[:500]}

--- 请将翻译结果替换下方 ---
[翻译结果]


"""
        return {"text": template, "template_mode": True}


class DocumentTranslator:
    """文档翻译器"""
    
    def __init__(self, engine_method: str = "auto"):
        self.engine = TranslationEngine(method=engine_method)
        self.hw = get_hardware_info()
        # 分段配置
        self.chunk_size = 2000 if self.hw.get("level") == "low" else 3000
    
    def detect_language(self, text: str) -> str:
        """简单语言检测"""
        if not text:
            return "unknown"
        
        # 统计中文字符比例
        chinese_chars = sum(1 for c in text if '一' <= c <= '鿿')
        ratio = chinese_chars / len(text) if text else 0
        
        if ratio > 0.3:
            return "zh"
        else:
            return "en"
    
    def translate_document(self, input_path: str, output_path: str,
                          src: str = "", tgt: str = "zh") -> Dict:
        """翻译文档"""
        input_path = str(Path(input_path).resolve())
        output_path = str(Path(output_path).resolve())
        
        if not os.path.exists(input_path):
            return {"success": False, "error": f"文件不存在: {input_path}"}
        
        ext = Path(input_path).suffix.lower()
        
        if ext == ".docx":
            return self._translate_docx(input_path, output_path, src, tgt)
        elif ext == ".xlsx":
            return self._translate_xlsx(input_path, output_path, src, tgt)
        elif ext == ".pptx":
            return self._translate_pptx(input_path, output_path, src, tgt)
        else:
            return {"success": False, "error": f"不支持的格式: {ext}"}
    
    def _translate_docx(self, input_path: str, output_path: str,
                       src: str, tgt: str) -> Dict:
        """翻译 Word 文档"""
        try:
            import docx
        except ImportError:
            return {"success": False, "error": "python-docx 未安装"}
        
        try:
            doc = docx.Document(input_path)
            
            # 翻译段落
            for para in doc.paragraphs:
                if para.text.strip():
                    # 检测语言
                    if not src:
                        src = self.detect_language(para.text)
                    
                    # 翻译
                    result = self.engine.translate(para.text, src, tgt)
                    if result.get("success") and not result.get("template_mode"):
                        # 保留格式，仅替换文本
                        if para.runs:
                            para.runs[0].text = result["text"]
                            for run in para.runs[1:]:
                                run.text = ""
                        else:
                            para.text = result["text"]
            
            # 翻译表格
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            result = self.engine.translate(cell.text, src, tgt)
                            if result.get("success") and not result.get("template_mode"):
                                cell.text = result["text"]
            
            doc.save(output_path)
            
            return {
                "success": True,
                "method": self.engine.get_best_method(),
                "output": output_path,
                "paragraphs_translated": len(doc.paragraphs),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _translate_xlsx(self, input_path: str, output_path: str,
                       src: str, tgt: str) -> Dict:
        """翻译 Excel 文档"""
        try:
            import openpyxl
        except ImportError:
            return {"success": False, "error": "openpyxl 未安装"}
        
        try:
            wb = openpyxl.load_workbook(input_path)
            translated_cells = 0
            
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            if not src:
                                src = self.detect_language(cell.value)
                            
                            result = self.engine.translate(cell.value, src, tgt)
                            if result.get("success") and not result.get("template_mode"):
                                cell.value = result["text"]
                                translated_cells += 1
            
            wb.save(output_path)
            
            return {
                "success": True,
                "method": self.engine.get_best_method(),
                "output": output_path,
                "cells_translated": translated_cells,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _translate_pptx(self, input_path: str, output_path: str,
                       src: str, tgt: str) -> Dict:
        """翻译 PPT 文档"""
        try:
            import pptx
        except ImportError:
            return {"success": False, "error": "python-pptx 未安装"}
        
        try:
            prs = pptx.Presentation(input_path)
            translated_shapes = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                if not src:
                                    src = self.detect_language(para.text)
                                
                                result = self.engine.translate(para.text, src, tgt)
                                if result.get("success") and not result.get("template_mode"):
                                    if para.runs:
                                        para.runs[0].text = result["text"]
                                        for run in para.runs[1:]:
                                            run.text = ""
                                    else:
                                        para.text = result["text"]
                                    translated_shapes += 1
            
            prs.save(output_path)
            
            return {
                "success": True,
                "method": self.engine.get_best_method(),
                "output": output_path,
                "shapes_translated": translated_shapes,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def batch_translate(self, input_dir: str, output_dir: str,
                       src: str = "", tgt: str = "zh") -> Dict:
        """批量翻译"""
        input_dir = Path(input_dir)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        supported = {".docx", ".xlsx", ".pptx"}
        files = [f for f in input_dir.iterdir() if f.suffix.lower() in supported]
        
        results = []
        for f in files:
            output_path = str(output_dir / f"{f.stem}_translated{f.suffix}")
            result = self.translate_document(str(f), output_path, src, tgt)
            results.append({"file": f.name, **result})
        
        success = sum(1 for r in results if r.get("success"))
        return {
            "success": True,
            "total": len(results),
            "success_count": success,
            "results": results,
        }


def _cli():
    """CLI 入口"""
    import argparse
    
    parser = argparse.ArgumentParser(description="文档翻译模块 v4.6.1")
    sub = parser.add_subparsers(dest="command", required=True)
    
    # translate
    p = sub.add_parser("translate", help="翻译文档")
    p.add_argument("--file", required=True, help="输入文件路径")
    p.add_argument("--output", required=True, help="输出文件路径")
    p.add_argument("--source", default="", help="源语言（不指定则自动检测）")
    p.add_argument("--target", default="zh", help="目标语言")
    p.add_argument("--method", default="auto", help="翻译引擎")
    
    # batch
    p = sub.add_parser("batch", help="批量翻译")
    p.add_argument("--input-dir", required=True)
    p.add_argument("--output-dir", required=True)
    p.add_argument("--source", default="")
    p.add_argument("--target", default="zh")
    p.add_argument("--method", default="auto")
    
    # detect
    p = sub.add_parser("detect", help="检测语言")
    p.add_argument("--file", required=True)
    
    # check
    p = sub.add_parser("check", help="检查可用引擎")
    
    args = parser.parse_args()
    
    if args.command == "translate":
        t = DocumentTranslator(engine_method=args.method)
        result = t.translate_document(args.file, args.output, args.source, args.target)
        print(json.dumps(result, ensure_ascii=False, default=str))
    
    elif args.command == "batch":
        t = DocumentTranslator(engine_method=args.method)
        result = t.batch_translate(args.input_dir, args.output_dir, args.source, args.target)
        print(json.dumps(result, ensure_ascii=False, default=str))
    
    elif args.command == "detect":
        t = DocumentTranslator()
        text = Path(args.file).read_text(encoding="utf-8")[:1000]
        lang = t.detect_language(text)
        print(json.dumps({"language": lang, "file": args.file}, ensure_ascii=False))
    
    elif args.command == "check":
        e = TranslationEngine()
        print(json.dumps({
            "available": e._available,
            "best": e.get_best_method(),
            "glossary_size": len(e.glossary),
        }, ensure_ascii=False))


if __name__ == "__main__":
    _cli()
